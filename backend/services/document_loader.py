"""
Document loader service for extracting text from various file formats
"""
import os
from typing import List, Dict, Any
from pathlib import Path
import docx
import pandas as pd
from pypdf import PdfReader
from pptx import Presentation

class DocumentLoader:
    """Service for loading and extracting text from documents"""
    
    @staticmethod
    def load_pdf(file_path: str) -> List[Dict[str, Any]]:
        """
        Extract text from PDF file
        
        Returns list of dicts with page_no and text
        """
        pages = []
        
        try:
            reader = PdfReader(file_path)
            
            for page_num, page in enumerate(reader.pages, start=1):
                text = page.extract_text()
                
                if text.strip():
                    pages.append({
                        "page_no": page_num,
                        "text": text,
                        "metadata": {
                            "source": os.path.basename(file_path),
                            "page": page_num,
                            "total_pages": len(reader.pages)
                        }
                    })
            
            return pages
        
        except Exception as e:
            raise Exception(f"Error loading PDF: {str(e)}")
    
    @staticmethod
    def load_docx(file_path: str) -> List[Dict[str, Any]]:
        """
        Extract text from DOCX file
        
        Returns list with single dict containing all text
        """
        try:
            doc = docx.Document(file_path)
            
            # Extract paragraphs
            paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
            
            # Extract tables
            table_texts = []
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text for cell in row.cells)
                    if row_text.strip():
                        table_texts.append(row_text)
            
            # Combine all text
            full_text = "\n".join(paragraphs)
            if table_texts:
                full_text += "\n\nTables:\n" + "\n".join(table_texts)
            
            return [{
                "page_no": 1,
                "text": full_text,
                "metadata": {
                    "source": os.path.basename(file_path),
                    "paragraphs": len(paragraphs),
                    "tables": len(doc.tables)
                }
            }]
        
        except Exception as e:
            raise Exception(f"Error loading DOCX: {str(e)}")
    
    @staticmethod
    def load_xlsx(file_path: str) -> List[Dict[str, Any]]:
        """
        Extract text from XLSX file
        
        Returns list with text from each sheet
        """
        try:
            excel_file = pd.ExcelFile(file_path)
            sheets_data = []
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                # Convert DataFrame to text representation
                text_lines = [f"Sheet: {sheet_name}\n"]
                
                # Add column headers
                text_lines.append(" | ".join(str(col) for col in df.columns))
                text_lines.append("-" * 50)
                
                # Add rows
                for _, row in df.iterrows():
                    row_text = " | ".join(str(val) for val in row.values)
                    text_lines.append(row_text)
                
                sheets_data.append({
                    "page_no": len(sheets_data) + 1,
                    "text": "\n".join(text_lines),
                    "metadata": {
                        "source": os.path.basename(file_path),
                        "sheet_name": sheet_name,
                        "rows": len(df),
                        "columns": len(df.columns)
                    }
                })
            
            return sheets_data
        
        except Exception as e:
            raise Exception(f"Error loading XLSX: {str(e)}")
    
    @staticmethod
    def load_pptx(file_path: str) -> List[Dict[str, Any]]:
        """
        Extract text from PPTX file
        
        Returns list with text from each slide
        """
        try:
            prs = Presentation(file_path)
            slides_data = []
            
            for slide_num, slide in enumerate(prs.slides, start=1):
                text_parts = []
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)
                
                if text_parts:
                    slides_data.append({
                        "page_no": slide_num,
                        "text": "\n".join(text_parts),
                        "metadata": {
                            "source": os.path.basename(file_path),
                            "slide_number": slide_num,
                            "total_slides": len(prs.slides)
                        }
                    })
            
            return slides_data
        
        except Exception as e:
            raise Exception(f"Error loading PPTX: {str(e)}")
    
    @classmethod
    def load_document(cls, file_path: str) -> List[Dict[str, Any]]:
        """
        Load document based on file extension
        
        Returns list of dicts with extracted text and metadata
        """
        file_ext = Path(file_path).suffix.lower()
        
        if file_ext == ".pdf":
            return cls.load_pdf(file_path)
        elif file_ext == ".docx":
            return cls.load_docx(file_path)
        elif file_ext in [".xlsx", ".xls"]:
            return cls.load_xlsx(file_path)
        elif file_ext == ".pptx":
            return cls.load_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file format: {file_ext}")
